#!/usr/bin/env python3
"""Generate china_trip.pptx — 4 slides about China education & AI / MAIC."""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import math

# ── constants ──────────────────────────────────────────────────────────────
SLIDE_W = Cm(33.87)
SLIDE_H = Cm(19.05)

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TITLE_CLR = RGBColor(0x1A, 0x1A, 0x18)
SECONDARY_CLR = RGBColor(0x73, 0x72, 0x6C)
TERTIARY_CLR = RGBColor(0x9C, 0x9A, 0x92)
CARD_BG = RGBColor(0xF1, 0xEF, 0xE8)

TAG_COLORS = {
    "blue":   {"bg": RGBColor(0xE6, 0xF1, 0xFB), "fg": RGBColor(0x18, 0x5F, 0xA5)},
    "purple": {"bg": RGBColor(0xEE, 0xED, 0xFE), "fg": RGBColor(0x53, 0x4A, 0xB7)},
    "green":  {"bg": RGBColor(0xE1, 0xF5, 0xEE), "fg": RGBColor(0x0F, 0x6E, 0x56)},
    "orange": {"bg": RGBColor(0xFA, 0xEE, 0xDA), "fg": RGBColor(0x85, 0x4F, 0x0B)},
    "grey":   {"bg": RGBColor(0xF1, 0xEF, 0xE8), "fg": RGBColor(0x5F, 0x5E, 0x5A)},
}

BAR_COLORS = {
    "blue":   RGBColor(0x37, 0x8A, 0xDD),
    "green":  RGBColor(0x5D, 0xCA, 0xA5),
    "orange": RGBColor(0xEF, 0x9F, 0x27),
    "red":    RGBColor(0xE2, 0x4B, 0x4A),
}

FONT = "Calibri"

# ── helpers ────────────────────────────────────────────────────────────────

def _set_font(run, size=11, bold=False, color=TITLE_CLR, name=FONT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name


def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None, radius=Cm(0.3)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    # adjust corner radius – adjustment[0] is corner size (0..100000)
    if shape.adjustments:
        # smaller value = less rounding; ~8000 gives a nice subtle radius
        shape.adjustments[0] = 0.06
    return shape


def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_paragraph(tf, text, size=11, bold=False, color=TITLE_CLR, align=PP_ALIGN.LEFT):
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    _set_font(r, size, bold, color)
    return p


def add_paragraph(tf, text, size=11, bold=False, color=TITLE_CLR, align=PP_ALIGN.LEFT, space_before=Pt(0)):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    r = p.add_run()
    r.text = text
    _set_font(r, size, bold, color)
    return p


def section_label(slide, left, top, width, text):
    """Small section header in uppercase grey."""
    tb = add_textbox(slide, left, top, width, Cm(0.6))
    set_paragraph(tb.text_frame, text, size=9, bold=True, color=SECONDARY_CLR)


def draw_card(slide, left, top, w, h, label, value, sub, border_color=None):
    """A metric card: small label, big number, subtitle."""
    add_rounded_rect(slide, left, top, w, h, CARD_BG, border_color=border_color)
    tb = add_textbox(slide, left + Cm(0.4), top + Cm(0.25), w - Cm(0.8), h - Cm(0.4))
    tf = tb.text_frame
    tf.word_wrap = True
    set_paragraph(tf, label, size=9, color=SECONDARY_CLR)
    add_paragraph(tf, value, size=20, bold=True, color=TITLE_CLR, space_before=Pt(2))
    add_paragraph(tf, sub, size=9, color=TERTIARY_CLR, space_before=Pt(1))


def draw_tag(slide, left, top, text, color_key):
    """Colored pill / tag."""
    colors = TAG_COLORS[color_key]
    tw = Cm(max(len(text) * 0.22 + 0.6, 1.6))
    th = Cm(0.55)
    shape = add_rounded_rect(slide, left, top, tw, th, colors["bg"])
    shape.text_frame.word_wrap = False
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = shape.text_frame.paragraphs[0].add_run()
    r.text = text
    _set_font(r, 8, True, colors["fg"])
    return tw


def draw_bar_row(slide, left, top, width, label, value, pct, bar_color, bar_frac):
    """Horizontal bar chart row: label | bar | percentage."""
    label_w = Cm(4.0)
    pct_w = Cm(2.0)
    bar_area_w = width - label_w - pct_w
    # label
    tb = add_textbox(slide, left, top, label_w, Cm(0.55))
    set_paragraph(tb.text_frame, label, size=9, color=TITLE_CLR, align=PP_ALIGN.RIGHT)
    # bar background
    bar_left = left + label_w + Cm(0.3)
    bar_top = top + Cm(0.05)
    bar_h = Cm(0.45)
    add_rounded_rect(slide, bar_left, bar_top, bar_area_w, bar_h, RGBColor(0xEC, 0xEB, 0xE4))
    # filled bar
    fill_w = int(bar_area_w * bar_frac)
    if fill_w > 0:
        shape = add_rounded_rect(slide, bar_left, bar_top, fill_w, bar_h, bar_color)
        # put value inside bar
        shape.text_frame.word_wrap = False
        shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        r = shape.text_frame.paragraphs[0].add_run()
        r.text = value
        _set_font(r, 8, True, WHITE)
    # percentage
    tb2 = add_textbox(slide, bar_left + bar_area_w + Cm(0.15), top, pct_w, Cm(0.55))
    set_paragraph(tb2.text_frame, pct, size=9, bold=True, color=SECONDARY_CLR)


def draw_text_card(slide, left, top, w, h, title_text, body_text, border_color=None):
    """Card with a title line and body paragraph."""
    add_rounded_rect(slide, left, top, w, h, CARD_BG, border_color=border_color)
    tb = add_textbox(slide, left + Cm(0.4), top + Cm(0.25), w - Cm(0.8), h - Cm(0.5))
    tf = tb.text_frame
    tf.word_wrap = True
    set_paragraph(tf, title_text, size=11, bold=True, color=TITLE_CLR)
    add_paragraph(tf, body_text, size=9, color=SECONDARY_CLR, space_before=Pt(4))


def draw_pipeline_block(slide, left, top, w, h, text, fill_color):
    shape = add_rounded_rect(slide, left, top, w, h, fill_color)
    shape.text_frame.word_wrap = True
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = shape.text_frame.paragraphs[0].add_run()
    r.text = text
    _set_font(r, 9, True, TITLE_CLR)


def draw_arrow(slide, left, top, w, h):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = TERTIARY_CLR
    shape.line.fill.background()


# ── slide builders ─────────────────────────────────────────────────────────

def build_slide1(prs):
    """Slide 1: China education & AI overview."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    LM = Cm(1.5)  # left margin
    CW = SLIDE_W - Cm(3.0)  # content width

    # Title
    tb = add_textbox(slide, LM, Cm(0.6), CW, Cm(1.0))
    set_paragraph(tb.text_frame, "Китай: образование и ИИ", size=24, bold=True, color=TITLE_CLR)
    tb2 = add_textbox(slide, LM, Cm(1.5), CW, Cm(0.5))
    set_paragraph(tb2.text_frame, "Ключевые цифры  \u00b7  данные Минобразования КНР, 2023\u20132025", size=12, color=SECONDARY_CLR)

    # ── Section: Population ──
    y = Cm(2.4)
    section_label(slide, LM, y, CW, "НАСЕЛЕНИЕ И СИСТЕМА ОБРАЗОВАНИЯ")
    y += Cm(0.7)
    card_w = (CW - Cm(0.9)) / 4
    card_h = Cm(2.2)
    cards = [
        ("Население Китая", "1.4 млрд", "2024"),
        ("Всего учащихся", "~340 млн", "все уровни, 2023\u20132024"),
        ("K-12 (1\u201312 класс)", "~195 млн", "начальная + средняя школа"),
        ("Дошкольники", "41 млн", "детские сады, 2023"),
    ]
    for i, (lbl, val, sub) in enumerate(cards):
        draw_card(slide, LM + i * (card_w + Cm(0.3)), y, card_w, card_h, lbl, val, sub)

    # ── Section: Funnel ──
    y += card_h + Cm(0.4)
    section_label(slide, LM, y, CW, "ВОРОНКА — ОТ ШКОЛЫ ДО ВУЗА")
    y += Cm(0.65)
    bars = [
        ("Сдают гаокао", "13.42 млн", "100%", BAR_COLORS["blue"], 1.0),
        ("Поступают в вузы", "~12 млн", "≈90%", BAR_COLORS["green"], 0.90),
        ("В топ-100 вузов", "~670 тыс.", "~5%", BAR_COLORS["orange"], 0.30),
        ("В топ 9 вузов C9", "", "<1%", BAR_COLORS["red"], 0.08),
    ]
    for label_t, val_t, pct_t, clr, frac in bars:
        draw_bar_row(slide, LM, y, CW, label_t, val_t, pct_t, clr, frac)
        y += Cm(0.65)

    # ── Section: Teachers & Unis ──
    y += Cm(0.15)
    section_label(slide, LM, y, CW, "УЧИТЕЛЯ И ВУЗЫ")
    y += Cm(0.65)
    card_h2 = Cm(2.2)
    cards2 = [
        ("Учителей всего", "18.9 млн", "на всех уровнях, 2023"),
        ("В школах K-12", "~16.8 млн", "нач. + ср. + ст. школа"),
        ("В вузах", "2.07 млн", "штатных преподавателей"),
        ("Всего вузов", "3 167", "по состоянию на 2025"),
    ]
    for i, (lbl, val, sub) in enumerate(cards2):
        draw_card(slide, LM + i * (card_w + Cm(0.3)), y, card_w, card_h2, lbl, val, sub)

    # ── Section: AI Companies ──
    y += card_h2 + Cm(0.4)
    section_label(slide, LM, y, CW, "КЛЮЧЕВЫЕ ИИ-КОМПАНИИ КИТАЯ")
    y += Cm(0.6)
    tb3 = add_textbox(slide, LM, y, CW, Cm(0.8))
    set_paragraph(
        tb3.text_frame,
        "Baidu, Alibaba, Tencent, ByteDance, Zhipu AI, Moonshot AI, MiniMax, Baichuan, 01.AI, SenseTime, iFlytek и ещё несколько десятков стартапов.",
        size=10, color=SECONDARY_CLR,
    )


def build_slide2(prs):
    """Slide 2: Where we visited."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    LM = Cm(1.5)
    CW = SLIDE_W - Cm(3.0)

    # Title
    tb = add_textbox(slide, LM, Cm(0.5), CW, Cm(1.0))
    set_paragraph(tb.text_frame, "Где побывали", size=24, bold=True, color=TITLE_CLR)
    tb2 = add_textbox(slide, LM, Cm(1.3), CW, Cm(0.5))
    set_paragraph(tb2.text_frame, "Китай  \u00b7  15\u201320 марта 2026  \u00b7  Шанхай, Ханчжоу, Шэньчжэнь, Пекин", size=11, color=SECONDARY_CLR)

    items = [
        ("ШАНХАЙ", [
            ("EdTech", "green",  "Squirrel AI", "Адаптивное ИИ-обучение K-12, китайский пионер персонализированного обучения"),
            ("Вузы",   "blue",   "CCN (Консорциум вузов)", "Посредник между китайскими и зарубежными университетами: двойные дипломы, обмен студентами"),
            ("ИИ",     "purple", "SAIRI", "Лаборатория ИИ при университете Цзяотун — компьютерное зрение и умное производство"),
            ("Вузы",   "blue",   "Shanghai Jiao Tong University", "Топовый исследовательский университет, ИИ для авиации, цифровые двойники, спутники Бэйдоу"),
            ("Школа",  "orange", "Luwan Middle + Senior School", "Школа 6\u201312 классов с упором на PBL, робототехнику и ИИ"),
        ]),
        ("ХАНЧЖОУ", [
            ("Вузы",   "blue",   "Hangzhou Normal University", "Педагогический университет, альма-матер Джека Ма"),
            ("Школа",  "orange", "Школа в Ханчжоу", "Школа под покровительством Джека Ма"),
        ]),
        ("ШЭНЬЧЖЭНЬ", [
            ("Вузы",   "blue",   "SGIS Tsinghua", "Международная аспирантура Цинхуа, 6 инженерных направлений под нужды индустрии Шэньчжэня"),
        ]),
        ("ПЕКИН", [
            ("Вузы",   "blue",   "Tsinghua University", "№1 в Китае, «китайский MIT», встреча с проф. Чжан — руководителем лаборатории ИИ в образовании"),
            ("ИИ",     "purple", "ByteDance", "ИИ-лаборатория, модели Doubao и Seedance, полный стек от инфраструктуры до продуктов"),
            ("EdTech", "green",  "Yuanfudao", "EdTech, онлайн-репетиторство K-12 и ИИ-помощник по домашним заданиям"),
        ]),
    ]

    y = Cm(2.1)
    for city, entries in items:
        section_label(slide, LM, y, CW, city)
        y += Cm(0.55)
        for tag_text, tag_color, name, desc in entries:
            tag_w = draw_tag(slide, LM, y, tag_text, tag_color)
            tb = add_textbox(slide, LM + tag_w + Cm(0.25), y - Cm(0.03), CW - tag_w - Cm(0.3), Cm(0.55))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r1 = p.add_run()
            r1.text = name + "  "
            _set_font(r1, 10, True, TITLE_CLR)
            r2 = p.add_run()
            r2.text = desc
            _set_font(r2, 9, False, SECONDARY_CLR)
            y += Cm(0.6)
        y += Cm(0.15)


def build_slide3(prs):
    """Slide 3: MAIC — what & how."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    LM = Cm(1.5)
    CW = SLIDE_W - Cm(3.0)

    # Title
    tb = add_textbox(slide, LM, Cm(0.5), CW, Cm(1.0))
    set_paragraph(tb.text_frame, "MAIC — Multi-Agent Interactive Classroom", size=22, bold=True, color=TITLE_CLR)
    tb2 = add_textbox(slide, LM, Cm(1.3), CW, Cm(0.5))
    set_paragraph(tb2.text_frame, "Разработка School of Education, Tsinghua University  \u00b7  2024\u20132025", size=11, color=SECONDARY_CLR)

    # ── Section: MOOC vs MAIC ──
    y = Cm(2.2)
    section_label(slide, LM, y, CW, "ИДЕЯ — ОТ ПАССИВНОГО ВИДЕО К ЖИВОМУ КЛАССУ")
    y += Cm(0.65)
    half_w = (CW - Cm(0.5)) / 2
    card_h = Cm(2.8)
    draw_text_card(
        slide, LM, y, half_w, card_h,
        "MOOC — было",
        "Студент смотрит записанное видео. Нет обратной связи, нет дискуссии, нет адаптации. Досматривают менее 10% записавшихся.",
    )
    draw_text_card(
        slide, LM + half_w + Cm(0.5), y, half_w, card_h,
        "MAIC — стало",
        "Полноценный ИИ-класс: учитель, ассистент и одноклассники — все агенты. Дискуссии, квизы, симуляции в реальном времени. Адаптируется под каждого студента.",
        border_color=RGBColor(0x7F, 0x77, 0xDD),
    )

    # ── Section: Agent roles ──
    y += card_h + Cm(0.35)
    section_label(slide, LM, y, CW, "РОЛИ АГЕНТОВ В КЛАССЕ")
    y += Cm(0.6)
    q_w = (CW - Cm(0.9)) / 4
    role_h = Cm(2.6)
    roles = [
        ("blue",   "Учитель",     "Ведёт урок, объясняет концепции, адаптирует стиль под аудиторию"),
        ("green",  "Ассистент",   "Помогает отстающим, проверяет задания, даёт мгновенную обратную связь"),
        ("purple", "Одноклассники", "4 типа: Sparker, Questioner, Thinker, Note-taker — дебаты и дискуссии"),
        ("orange", "Recall Agent", "Возвращает отключившихся студентов: +79% повторных логинов"),
    ]
    for i, (clr_key, title_t, body_t) in enumerate(roles):
        x = LM + i * (q_w + Cm(0.3))
        add_rounded_rect(slide, x, y, q_w, role_h, CARD_BG)
        # tag at top of card
        draw_tag(slide, x + Cm(0.3), y + Cm(0.25), title_t, clr_key)
        tb = add_textbox(slide, x + Cm(0.3), y + Cm(0.95), q_w - Cm(0.6), role_h - Cm(1.1))
        tf = tb.text_frame
        tf.word_wrap = True
        set_paragraph(tf, body_t, size=9, color=SECONDARY_CLR)

    # ── Section: Capabilities ──
    y += role_h + Cm(0.35)
    section_label(slide, LM, y, CW, "ЧТО УМЕЕТ ПЛАТФОРМА")
    y += Cm(0.6)
    third_w = (CW - Cm(0.6)) / 3
    cap_h = Cm(3.0)
    caps = [
        ("Генерация курса", "Загрузи PDF или опиши тему — через ~30 минут готов полный интерактивный урок со слайдами, голосом, квизами и симуляциями. Стоимость: менее $2 в API."),
        ("Адаптивность", "Движок на основе таксономии Блума и зоны ближайшего развития (ZPD). ИИ отслеживает уровень каждого студента и меняет сложность в реальном времени."),
        ("Открытый код", "Open-source, AGPL-3.0. Поддерживает OpenAI, Anthropic, DeepSeek, Gemini. Деплой через Docker или Vercel. Экспорт в PPTX и HTML."),
    ]
    for i, (t, b) in enumerate(caps):
        draw_text_card(slide, LM + i * (third_w + Cm(0.3)), y, third_w, cap_h, t, b)


def build_slide4(prs):
    """Slide 4: MAIC — results."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    LM = Cm(1.5)
    CW = SLIDE_W - Cm(3.0)

    # Title
    tb = add_textbox(slide, LM, Cm(0.5), CW, Cm(1.0))
    set_paragraph(tb.text_frame, "MAIC — результаты", size=24, bold=True, color=TITLE_CLR)
    tb2 = add_textbox(slide, LM, Cm(1.3), CW, Cm(0.5))
    set_paragraph(tb2.text_frame, "2+ года тестирования на студентах Tsinghua University", size=11, color=SECONDARY_CLR)

    # ── Pipeline ──
    y = Cm(2.2)
    section_label(slide, LM, y, CW, "КАК РАБОТАЕТ ПАЙПЛАЙН")
    y += Cm(0.6)
    block_w = Cm(6.5)
    arrow_w = Cm(1.0)
    block_h = Cm(1.3)
    arrow_h = Cm(0.5)
    labels = ["PDF / тема", "Lesson Planner\n(план урока)", "Scene Producer\n(мультимедиа)", "Интерактивный класс\n(28+ типов действий)"]
    colors_pipe = [
        TAG_COLORS["blue"]["bg"], TAG_COLORS["green"]["bg"],
        TAG_COLORS["purple"]["bg"], TAG_COLORS["orange"]["bg"],
    ]
    total_blocks_w = 4 * block_w + 3 * arrow_w
    start_x = LM + (CW - total_blocks_w) // 2
    for i, (lbl, clr) in enumerate(zip(labels, colors_pipe)):
        bx = start_x + i * (block_w + arrow_w)
        draw_pipeline_block(slide, bx, y, block_w, block_h, lbl, clr)
        if i < 3:
            ax = bx + block_w
            draw_arrow(slide, ax + Cm(0.1), y + (block_h - arrow_h) // 2, arrow_w - Cm(0.2), arrow_h)

    # ── Scale ──
    y += block_h + Cm(0.5)
    section_label(slide, LM, y, CW, "МАСШТАБ")
    y += Cm(0.6)
    card_w = (CW - Cm(0.9)) / 4
    card_h = Cm(2.1)
    scale_cards = [
        ("Студентов охвачено", "90 000+", "к декабрю 2025"),
        ("Учебных записей", "100 000+", "реальные взаимодействия"),
        ("Удовлетворённость", "84\u201392%", "среди активных пользователей"),
        ("GitHub звёзды", "13 300+", "за 2 недели после релиза"),
    ]
    for i, (lbl, val, sub) in enumerate(scale_cards):
        draw_card(slide, LM + i * (card_w + Cm(0.3)), y, card_w, card_h, lbl, val, sub)

    # ── Key Result ──
    y += card_h + Cm(0.4)
    section_label(slide, LM, y, CW, "КЛЮЧЕВОЙ РЕЗУЛЬТАТ")
    y += Cm(0.6)
    third_w = (CW - Cm(0.6)) / 3
    res_h = Cm(2.8)
    results = [
        ("grey",   "Традиционный MOOC", "Пассивный просмотр. Менее 10% досматривают до конца."),
        ("green",  "MAIC — ИИ-класс", "Значимо лучше MOOC. Сопоставимо с живым преподавателем при масштабировании на тысячи студентов."),
        ("purple", "Открытие", "Студенты с низким начальным уровнем опирались на совместное построение знаний с агентами — и показали наибольший прирост."),
    ]
    for i, (clr_key, title_t, body_t) in enumerate(results):
        x = LM + i * (third_w + Cm(0.3))
        border_clr = TAG_COLORS[clr_key]["fg"] if clr_key != "grey" else None
        add_rounded_rect(slide, x, y, third_w, res_h, CARD_BG, border_color=border_clr)
        draw_tag(slide, x + Cm(0.3), y + Cm(0.25), title_t, clr_key)
        tb = add_textbox(slide, x + Cm(0.3), y + Cm(0.95), third_w - Cm(0.6), res_h - Cm(1.1))
        tf = tb.text_frame
        tf.word_wrap = True
        set_paragraph(tf, body_t, size=9, color=SECONDARY_CLR)

    # ── Publications ──
    y += res_h + Cm(0.4)
    section_label(slide, LM, y, CW, "ПУБЛИКАЦИИ И ПРИЗНАНИЕ")
    y += Cm(0.6)
    half_w = (CW - Cm(0.3)) / 2
    pub_h = Cm(2.4)
    draw_text_card(
        slide, LM, y, half_w, pub_h,
        "Публикации",
        "From MOOC to MAIC — JCST 2026  /  SimClass — NAACL 2025  /  Slide2Lecture — KDD 2025  /  AI as Learning Partners — ICLS 2025",
    )
    draw_text_card(
        slide, LM + half_w + Cm(0.3), y, half_w, pub_h,
        "Государственное признание",
        "В марте 2025 MAIC вошла в число первых приложений на национальной платформе «умного образования» Китая (国家智慧教育公共服务平台).",
    )


# ── main ───────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_slide1(prs)
    build_slide2(prs)
    build_slide3(prs)
    build_slide4(prs)

    out = "/home/user/china_trip.pptx"
    prs.save(out)
    print(f"Saved to {out}")


if __name__ == "__main__":
    main()

#!/usr/bin/env python3
"""Generate squirrel_ai.pptx — slides about Squirrel AI Learning."""

from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

SLIDE_W = Cm(33.87)
SLIDE_H = Cm(19.05)

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TITLE_CLR = RGBColor(0x1A, 0x1A, 0x18)
SECONDARY_CLR = RGBColor(0x73, 0x72, 0x6C)
TERTIARY_CLR = RGBColor(0x9C, 0x9A, 0x92)
CARD_BG = RGBColor(0xF1, 0xEF, 0xE8)

TAG_COLORS = {
    "blue":   {"bg": RGBColor(0xE6, 0xF1, 0xFB), "fg": RGBColor(0x18, 0x5F, 0xA5)},
    "purple": {"bg": RGBColor(0xEE, 0xED, 0xFE), "fg": RGBColor(0x53, 0x4A, 0xB7)},
    "green":  {"bg": RGBColor(0xE1, 0xF5, 0xEE), "fg": RGBColor(0x0F, 0x6E, 0x56)},
    "orange": {"bg": RGBColor(0xFA, 0xEE, 0xDA), "fg": RGBColor(0x85, 0x4F, 0x0B)},
    "grey":   {"bg": RGBColor(0xF1, 0xEF, 0xE8), "fg": RGBColor(0x5F, 0x5E, 0x5A)},
    "red":    {"bg": RGBColor(0xFD, 0xE8, 0xE8), "fg": RGBColor(0xB9, 0x1C, 0x1C)},
}

FONT = "Calibri"


def _sf(run, size=11, bold=False, color=TITLE_CLR, name=FONT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name


def add_rr(slide, left, top, w, h, fill, border=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    if shape.adjustments:
        shape.adjustments[0] = 0.06
    return shape


def tb(slide, left, top, w, h):
    return slide.shapes.add_textbox(left, top, w, h)


def setp(tf, text, size=11, bold=False, color=TITLE_CLR, align=PP_ALIGN.LEFT):
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    _sf(r, size, bold, color)
    return p


def addp(tf, text, size=11, bold=False, color=TITLE_CLR, align=PP_ALIGN.LEFT, sb=Pt(0)):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = sb
    r = p.add_run()
    r.text = text
    _sf(r, size, bold, color)
    return p


def slabel(slide, left, top, w, text):
    t = tb(slide, left, top, w, Cm(0.6))
    setp(t.text_frame, text, size=9, bold=True, color=SECONDARY_CLR)


def card(slide, left, top, w, h, label, value, sub, border=None):
    add_rr(slide, left, top, w, h, CARD_BG, border)
    t = tb(slide, left + Cm(0.4), top + Cm(0.25), w - Cm(0.8), h - Cm(0.4))
    tf = t.text_frame
    tf.word_wrap = True
    setp(tf, label, size=9, color=SECONDARY_CLR)
    addp(tf, value, size=20, bold=True, color=TITLE_CLR, sb=Pt(2))
    addp(tf, sub, size=9, color=TERTIARY_CLR, sb=Pt(1))


def tag(slide, left, top, text, ck):
    c = TAG_COLORS[ck]
    tw = Cm(max(len(text) * 0.22 + 0.6, 1.6))
    th = Cm(0.55)
    s = add_rr(slide, left, top, tw, th, c["bg"])
    s.text_frame.word_wrap = False
    s.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = s.text_frame.paragraphs[0].add_run()
    r.text = text
    _sf(r, 8, True, c["fg"])
    return tw


def tcard(slide, left, top, w, h, title, body, border=None):
    add_rr(slide, left, top, w, h, CARD_BG, border)
    t = tb(slide, left + Cm(0.4), top + Cm(0.25), w - Cm(0.8), h - Cm(0.5))
    tf = t.text_frame
    tf.word_wrap = True
    setp(tf, title, size=11, bold=True, color=TITLE_CLR)
    addp(tf, body, size=9, color=SECONDARY_CLR, sb=Pt(4))


def tag_card(slide, left, top, w, h, tag_text, tag_color, body):
    add_rr(slide, left, top, w, h, CARD_BG)
    tag(slide, left + Cm(0.3), top + Cm(0.25), tag_text, tag_color)
    t = tb(slide, left + Cm(0.3), top + Cm(0.95), w - Cm(0.6), h - Cm(1.1))
    tf = t.text_frame
    tf.word_wrap = True
    setp(tf, body, size=9, color=SECONDARY_CLR)


def arrow(slide, left, top, w, h):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = TERTIARY_CLR
    s.line.fill.background()


def pipe_block(slide, left, top, w, h, text, fill):
    s = add_rr(slide, left, top, w, h, fill)
    s.text_frame.word_wrap = True
    s.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = s.text_frame.paragraphs[0].add_run()
    r.text = text
    _sf(r, 9, True, TITLE_CLR)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1: Overview
# ═══════════════════════════════════════════════════════════════════════════
def slide1(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = WHITE
    LM = Cm(1.5)
    CW = SLIDE_W - Cm(3.0)

    t = tb(s, LM, Cm(0.5), CW, Cm(1.0))
    setp(t.text_frame, "Squirrel AI — адаптивное обучение на базе ИИ", size=22, bold=True, color=TITLE_CLR)
    t2 = tb(s, LM, Cm(1.3), CW, Cm(0.5))
    setp(t2.text_frame, "松鼠AI  \u00b7  Шанхай, Китай  \u00b7  основана в 2014  \u00b7  первый EdTech-единорог Китая в адаптивном обучении", size=11, color=SECONDARY_CLR)

    # Section: Key numbers
    y = Cm(2.2)
    slabel(s, LM, y, CW, "КОМПАНИЯ В ЦИФРАХ")
    y += Cm(0.65)
    cw = (CW - Cm(1.2)) / 5
    ch = Cm(2.2)
    nums = [
        ("Студентов", "43 млн", "по всему миру"),
        ("Школ-партнёров", "60 000", "используют платформу"),
        ("Учебных центров", "1 900+", "в 300+ городах Китая"),
        ("Привлечено", "$128 млн", "Series C-II"),
        ("Данных обучения", "10 млрд", "поведенческих записей"),
    ]
    for i, (l, v, su) in enumerate(nums):
        card(s, LM + i * (cw + Cm(0.3)), y, cw, ch, l, v, su)

    # Section: What it does
    y += ch + Cm(0.4)
    slabel(s, LM, y, CW, "ЧТО ДЕЛАЕТ SQUIRREL AI")
    y += Cm(0.65)
    tw = (CW - Cm(0.6)) / 3
    th = Cm(3.2)
    items = [
        ("Адаптивное K-12 обучение",
         "ИИ-платформа для персонализированного обучения школьников. Разбивает предметы на тысячи «нано-знаний» (10 000+ точек только по математике) и строит индивидуальный маршрут для каждого ученика."),
        ("Диагностика и рекомендации",
         "Система находит конкретные пробелы в знаниях ученика и выстраивает оптимальный путь их устранения. Точность подбора заданий выросла с 78% до 93% благодаря Large Adaptive Model."),
        ("Планшет + ИИ-репетитор",
         "Ученик работает на планшете Squirrel AI. Система отслеживает эмоции (скука, фрустрация), принимает голосовой и текстовый ввод, ведёт 100+ сценариев диалога."),
    ]
    for i, (ti, bo) in enumerate(items):
        tcard(s, LM + i * (tw + Cm(0.3)), y, tw, th, ti, bo)

    # Section: Recognition
    y += th + Cm(0.4)
    slabel(s, LM, y, CW, "ПРИЗНАНИЕ")
    y += Cm(0.6)
    t3 = tb(s, LM, y, CW, Cm(1.2))
    tf = t3.text_frame
    tf.word_wrap = True
    setp(tf, "TIME Best Inventions 2025  \u00b7  5 статей на ACL 2025  \u00b7  IEEE P3428 Standard (председатель — основатель Squirrel AI)  \u00b7  Совместная лаборатория с CMU  \u00b7  Кейс Stanford GSB  \u00b7  Выход на рынок США запланирован на 2026", size=10, color=SECONDARY_CLR)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2: Architecture (based on photos)
# ═══════════════════════════════════════════════════════════════════════════
def slide2(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = WHITE
    LM = Cm(1.5)
    CW = SLIDE_W - Cm(3.0)

    t = tb(s, LM, Cm(0.5), CW, Cm(1.0))
    setp(t.text_frame, "Squirrel AI Engine — архитектура платформы", size=22, bold=True, color=TITLE_CLR)
    t2 = tb(s, LM, Cm(1.3), CW, Cm(0.5))
    setp(t2.text_frame, "По материалам презентации в офисе Squirrel AI, Шанхай", size=11, color=SECONDARY_CLR)

    # Layer labels
    y = Cm(2.2)
    slabel(s, LM, y, CW, "ТРИ СЛОЯ ПЛАТФОРМЫ")
    y += Cm(0.65)

    # Top layer: Knowledge & Content
    third_w = (CW - Cm(0.6)) / 3
    layer_h = Cm(3.6)

    tcard(s, LM, y, third_w, layer_h,
        "Слой знаний (Knowledge Layer)",
        "Knowledge Graph: 10 000+ точек знаний по математике, связанных в граф зависимостей. "
        "Content Map (Item Bank), Learning Map, Mistake Reasoning Ontology — онтологии контента, "
        "учебных путей и типичных ошибок.",
        border=TAG_COLORS["blue"]["fg"])

    tcard(s, LM + third_w + Cm(0.3), y, third_w, layer_h,
        "Слой моделей (Model Layer)",
        "Large Adaptive Model (LAM) — первая в мире мультимодальная адаптивная модель для образования. "
        "Включает: Knowledge Tracing (Learner Model), Recommendation Engine, "
        "Assessment Engine, Teaching Routing Engine, RAG-пайплайн с MoE multi-modal LLMs.",
        border=TAG_COLORS["purple"]["fg"])

    tcard(s, LM + 2 * (third_w + Cm(0.3)), y, third_w, layer_h,
        "Слой агентов (Agent Layer)",
        "Squirrel AI Adaptive Agents: Data Analysis Agent, Teaching Autonomous Driving Agent, "
        "Mathematics Q&A Agent, Chinese Language Comprehension Agent. Коллаборативные ИИ-агенты "
        "для аналитики, преподавания, рассуждений и проверки понимания.",
        border=TAG_COLORS["green"]["fg"])

    # Pipeline
    y += layer_h + Cm(0.5)
    slabel(s, LM, y, CW, "КАК РАБОТАЕТ АДАПТИВНЫЙ ЦИКЛ")
    y += Cm(0.65)

    bw = Cm(5.2)
    aw = Cm(0.9)
    bh = Cm(1.3)
    ah = Cm(0.45)
    labels = [
        "Диагностика\n(Assessment)",
        "Профиль ученика\n(Knowledge Tracing)",
        "Рекомендация\n(Learning Path)",
        "Обучение\n(Adaptive Content)",
        "Обратная связь\n(Feedback Loop)",
    ]
    fills = [
        TAG_COLORS["blue"]["bg"], TAG_COLORS["purple"]["bg"],
        TAG_COLORS["green"]["bg"], TAG_COLORS["orange"]["bg"],
        TAG_COLORS["blue"]["bg"],
    ]
    total = 5 * bw + 4 * aw
    sx = LM + (CW - total) // 2
    for i, (lb, fl) in enumerate(zip(labels, fills)):
        bx = sx + i * (bw + aw)
        pipe_block(s, bx, y, bw, bh, lb, fl)
        if i < 4:
            arrow(s, bx + bw + Cm(0.05), y + (bh - ah) // 2, aw - Cm(0.1), ah)

    # Bottom: LMS Core
    y += bh + Cm(0.5)
    slabel(s, LM, y, CW, "ИНФРАСТРУКТУРА (LMS CORE)")
    y += Cm(0.6)
    t3 = tb(s, LM, y, CW, Cm(1.0))
    tf = t3.text_frame
    tf.word_wrap = True
    setp(tf,
        "LRS (Learning Record Store)  \u00b7  Real-time Event Collector  \u00b7  MIBA Event Collector  \u00b7  "
        "Dialog-Based HUI  \u00b7  Assessment Service  \u00b7  Navigation Service  \u00b7  "
        "Monitoring & Alert Services  \u00b7  LMS Presentation & Administration Services  \u00b7  "
        "Control & Feedback System  \u00b7  Supervision System  \u00b7  Human-in-the-Loop Validation",
        size=9, color=SECONDARY_CLR)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3: Results (based on tablet photos)
# ═══════════════════════════════════════════════════════════════════════════
def slide3(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = WHITE
    LM = Cm(1.5)
    CW = SLIDE_W - Cm(3.0)

    t = tb(s, LM, Cm(0.5), CW, Cm(1.0))
    setp(t.text_frame, "Squirrel AI — результаты адаптивного обучения", size=22, bold=True, color=TITLE_CLR)
    t2 = tb(s, LM, Cm(1.3), CW, Cm(0.5))
    setp(t2.text_frame, "Данные с планшета ученика, демонстрация в офисе Squirrel AI, Шанхай", size=11, color=SECONDARY_CLR)

    # Before/After table
    y = Cm(2.2)
    slabel(s, LM, y, CW, "ДО И ПОСЛЕ: ПРОГРЕСС ПО УЧЕБНЫМ ЦЕЛЯМ (LEARNING OBJECTIVES)")
    y += Cm(0.7)

    # Table-like layout with cards
    half_w = (CW - Cm(0.5)) / 2

    # Before card
    add_rr(s, LM, y, half_w, Cm(5.5), CARD_BG, border=TAG_COLORS["red"]["fg"])
    tag(s, LM + Cm(0.3), y + Cm(0.25), "До обучения (Before)", "red")
    items_before = [
        "Compare two quantities: Progress 70%, Proficiency: Average",
        "Compare quantities to find equality: Progress 94%, Proficiency: Good",
        "Compare two numbers using words: Progress 35%, Proficiency: Unmastered",
        "Introduction to comparison symbols: Progress 35%, Proficiency: Unmastered",
        "Compare numbers using symbols: Progress 35%, Proficiency: Unmastered",
        "Fill in missing symbols/numbers: Progress 35%, Proficiency: Unmastered",
    ]
    ty = y + Cm(1.0)
    for item in items_before:
        t3 = tb(s, LM + Cm(0.4), ty, half_w - Cm(0.8), Cm(0.55))
        setp(t3.text_frame, "\u2022 " + item, size=8, color=SECONDARY_CLR)
        ty += Cm(0.55)

    # Summary line before
    ty += Cm(0.2)
    t4 = tb(s, LM + Cm(0.4), ty, half_w - Cm(0.8), Cm(0.5))
    setp(t4.text_frame, "Средний прогресс: 35\u201370%  \u00b7  Большинство тем: Unmastered", size=9, bold=True, color=TAG_COLORS["red"]["fg"])

    # After card
    add_rr(s, LM + half_w + Cm(0.5), y, half_w, Cm(5.5), CARD_BG, border=TAG_COLORS["green"]["fg"])
    tag(s, LM + half_w + Cm(0.8), y + Cm(0.25), "После обучения (After)", "green")
    items_after = [
        "Compare two quantities: Progress 94%, Proficiency: Good",
        "Compare quantities to find equality: \u2014 (уже освоено)",
        "Compare two numbers using words: Progress 91%, Proficiency: Good",
        "Introduction to comparison symbols: Progress 94%, Proficiency: Good",
        "Compare numbers using symbols: Progress 91%, Proficiency: Good",
        "Fill in missing symbols/numbers: Progress 91%, Proficiency: Good",
    ]
    ty = y + Cm(1.0)
    ax = LM + half_w + Cm(0.5)
    for item in items_after:
        t3 = tb(s, ax + Cm(0.4), ty, half_w - Cm(0.8), Cm(0.55))
        setp(t3.text_frame, "\u2022 " + item, size=8, color=SECONDARY_CLR)
        ty += Cm(0.55)

    ty += Cm(0.2)
    t4 = tb(s, ax + Cm(0.4), ty, half_w - Cm(0.8), Cm(0.5))
    setp(t4.text_frame, "Средний прогресс: 91\u201394%  \u00b7  Все темы: Good / Mastered", size=9, bold=True, color=TAG_COLORS["green"]["fg"])

    # Improvement cards
    y += Cm(6.2)
    slabel(s, LM, y, CW, "УЛУЧШЕНИЕ ПО ТЕМАМ (IMPROVEMENT)")
    y += Cm(0.65)
    cw4 = (CW - Cm(0.9)) / 4
    ch2 = Cm(2.1)
    improvements = [
        ("Compare quantities", "+24%", "с 70% до 94%"),
        ("Comparison using words", "+56%", "с 35% до 91%"),
        ("Comparison symbols", "+59%", "с 35% до 94%"),
        ("Numbers using symbols", "+56%", "с 35% до 91%"),
    ]
    for i, (l, v, su) in enumerate(improvements):
        card(s, LM + i * (cw4 + Cm(0.3)), y, cw4, ch2, l, v, su, border=TAG_COLORS["green"]["fg"])

    # Answer accuracy
    y += ch2 + Cm(0.4)
    slabel(s, LM, y, CW, "ТОЧНОСТЬ ОТВЕТОВ (ANSWER ACCURACY)")
    y += Cm(0.6)
    t5 = tb(s, LM, y, CW, Cm(0.6))
    setp(t5.text_frame,
        "После адаптивного обучения: 88\u2013100% точности ответов по всем освоенным темам  \u00b7  "
        "Время на тему: от 16 секунд до 2 мин 19 сек",
        size=10, color=SECONDARY_CLR)


# ═══════════════════════════════════════════════════════════════════════════
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide1(prs)
    slide2(prs)
    slide3(prs)
    out = os.path.join(os.path.expanduser("~"), "squirrel_ai.pptx")
    prs.save(out)
    print(f"Saved to {out}")


if __name__ == "__main__":
    main()
